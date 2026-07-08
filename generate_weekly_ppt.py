"""
周报 PPT 生成器
根据用户提供的周报数据，生成与现有模板风格一致的 3 页 PPTX 文件。
支持在条目中嵌入语雀文档链接，可在 PPT 中点击跳转。

用法:
  python generate_weekly_ppt.py --date-range "7.7-7.11" \
      --this-week "工作项1" "工作项2" \
      --next-week "计划项1" "计划项2" \
      --long-term "长期任务1" \
      --output "周报_7.7-7.11.pptx"

也可以从 JSON 文件读取:
  python generate_weekly_ppt.py --json data.json --output "周报.pptx"
"""

import argparse
import json
import sys
import io
from pathlib import Path

# 修复 Windows GBK 编码问题
if sys.stdout.encoding != 'utf-8':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ============================================================
# 样式常量（从 周报7.pptx 模板提取）
# ============================================================

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH  = 12192000  # EMU
SLIDE_HEIGHT = 6858000   # EMU

# 色彩方案 (WPS 主题)
COLOR_DK2     = RGBColor(0x44, 0x54, 0x6A)  # 深色文字主色
COLOR_ACCENT1 = RGBColor(0x48, 0x74, 0xCB)  # 蓝色强调
COLOR_ACCENT2 = RGBColor(0xEE, 0x82, 0x2F)  # 橙色强调
COLOR_LINK    = RGBColor(0x05, 0x63, 0xC1)  # 超链接蓝色
COLOR_SUB     = RGBColor(0x59, 0x56, 0x59)  # 子条目灰
COLOR_META    = RGBColor(0x8A, 0x8A, 0x8A)  # 元信息浅灰

# 标题样式
TITLE_FONT_NAME  = "微软雅黑"
TITLE_FONT_SIZE  = Pt(44)
TITLE_FONT_BOLD  = True
TITLE_FONT_COLOR = COLOR_DK2

# 内容样式
CONTENT_FONT_NAME = "微软雅黑"
CONTENT_FONT_SIZE = Pt(18)
CONTENT_FONT_BOLD = False
CONTENT_FONT_COLOR = COLOR_DK2

# 子条目样式
SUBITEM_FONT_SIZE = Pt(16)
SUBITEM_FONT_COLOR = COLOR_SUB

# 链接标签样式
LINK_LABEL_FONT_SIZE = Pt(13)
LINK_URL_FONT_SIZE   = Pt(10)

# 标题位置 (EMU)
TITLE_LEFT   = 1200000
TITLE_TOP    = 250000
TITLE_WIDTH  = 9800000
TITLE_HEIGHT = 770000

# 内容区域位置
CONTENT_LEFT   = 1200000
CONTENT_TOP_SLIDE1 = 1050000
CONTENT_WIDTH  = 9800000
CONTENT_HEIGHT_SLIDE1 = 5400000  # 第1页内容区较大
CONTENT_TOP_OTHER   = 1850000
CONTENT_HEIGHT_OTHER = 3200000  # 第2、3页内容区较小

# 行间距
LINE_SPACING = Pt(36)

# 最大条目数（超过后缩小字号）
MAX_ITEMS_BEFORE_SHRINK = 8


# ============================================================
# 工具函数
# ============================================================

def _add_textbox(slide, left, top, width, height):
    """在幻灯片上添加文本框"""
    return slide.shapes.add_textbox(left, top, width, height)


def _set_run(run, text, font_name, font_size, font_color, bold=False):
    """设置 run 的字体样式"""
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold


def _make_hyperlink_run(run, label, url):
    """
    将 run 设置为可点击的超链接。

    在 python-pptx 中，超链接通过操作底层 XML 添加:
      <a:rPr> 下添加 <a:hlinkClick r:id="..." />
    同时需要在 slide 的 relationship 中添加外部链接。
    但 python-pptx 的 run.hyperlink 属性直接支持 setter。
    """
    run.text = label
    run.font.name = "微软雅黑"
    run.font.size = LINK_LABEL_FONT_SIZE
    run.font.color.rgb = COLOR_LINK
    run.font.underline = True
    run.font.bold = False
    # 设置超链接
    run.hyperlink.address = url


def _add_paragraph_with_link(tf, prefix, text, link_url, font_size, is_first=False):
    """
    添加一个带链接的段落。

    格式: "prefixtext  📎语雀文档"
    其中 "📎语雀文档" 是可点击超链接。
    """
    if not is_first:
        para = tf.add_paragraph()
    else:
        para = tf.paragraphs[0]

    para.alignment = PP_ALIGN.LEFT
    para.space_after = Pt(4)

    # 正文 run
    body_run = para.add_run()
    body_run.text = f"{prefix}{text}"
    body_run.font.name = "微软雅黑"
    body_run.font.size = font_size
    body_run.font.color.rgb = COLOR_ACCENT1
    body_run.font.bold = True

    if link_url:
        # 添加一个空格
        spacer = para.add_run()
        spacer.text = "  "
        spacer.font.size = Pt(8)

        # 可点击的链接标签
        link_run = para.add_run()
        _make_hyperlink_run(link_run, "📎语雀文档", link_url)

    return para


def _add_sub_paragraph(tf, text, link_url, font_size, indent=True):
    """添加子条目段落，支持可选的超链接"""
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    para.space_after = Pt(2)

    prefix = "     " if indent else ""

    if isinstance(text, dict):
        text = text.get("text", "")
        link_url = text.get("link", link_url)

    # 子条目文本
    body_run = para.add_run()
    body_run.text = f"{prefix}{text}"
    body_run.font.name = "微软雅黑"
    body_run.font.size = font_size
    body_run.font.color.rgb = SUBITEM_FONT_COLOR
    body_run.font.bold = False

    if link_url:
        spacer = para.add_run()
        spacer.text = "  "
        spacer.font.size = Pt(6)

        link_run = para.add_run()
        _make_hyperlink_run(link_run, "📎链接", link_url)

    return para


def _add_link_only_line(tf, link_url, font_size):
    """当主条目没有 link 但想要附上链接时，添加一条独立的链接行"""
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    para.space_after = Pt(6)

    # URL 文本（灰显）
    url_run = para.add_run()
    url_run.text = f"     🔗 "
    url_run.font.name = "微软雅黑"
    url_run.font.size = LINK_URL_FONT_SIZE
    url_run.font.color.rgb = COLOR_META

    link_run = para.add_run()
    _make_hyperlink_run(link_run, link_url, link_url)


def _add_title_slide(slide, title_text, content_items, is_first_slide=False):
    """
    添加标题+内容幻灯片

    Args:
        slide: pptx Slide 对象
        title_text: 标题文字
        content_items: 内容列表，每项可以是:
            - str: 单行文本
            - dict: {
                "text": "主条目",
                "link": "https://...",    # 可选：语雀文档链接
                "sub": [
                    "子条目1",
                    {"text": "子条目2", "link": "https://..."}  # 子条目也可带链接
                ]
              }
        is_first_slide: 是否第1页（影响内容区域大小和布局）
    """
    # 添加标题
    title_box = _add_textbox(slide, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = title_box.text_frame
    tf.word_wrap = True
    title_para = tf.paragraphs[0]
    _set_run(title_para.add_run(), title_text,
             TITLE_FONT_NAME, TITLE_FONT_SIZE, TITLE_FONT_COLOR, bold=True)
    title_para.alignment = PP_ALIGN.LEFT

    # 添加内容
    content_top = CONTENT_TOP_SLIDE1 if is_first_slide else CONTENT_TOP_OTHER
    content_height = CONTENT_HEIGHT_SLIDE1 if is_first_slide else CONTENT_HEIGHT_OTHER
    content_box = _add_textbox(slide, CONTENT_LEFT, content_top,
                               CONTENT_WIDTH, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    # 决定字号（内容多时自动缩小）
    total_lines = sum(
        len(item.get("sub", [])) + 1 + (1 if isinstance(item, dict) and item.get("link") else 0)
        if isinstance(item, dict) else 1
        for item in content_items
    )
    shrink = total_lines > MAX_ITEMS_BEFORE_SHRINK
    item_font_size = Pt(16) if shrink else CONTENT_FONT_SIZE
    sub_font_size = Pt(14) if shrink else SUBITEM_FONT_SIZE

    first_para = True

    for idx, item in enumerate(content_items):
        if isinstance(item, str):
            item = {"text": item, "sub": []}

        text = item.get("text", "")
        link = item.get("link", "")
        subs = item.get("sub", [])

        # 主条目（带链接）
        _add_paragraph_with_link(tf, "▎", text, link, item_font_size,
                                 is_first=first_para)
        first_para = False

        # 子条目
        for sub in subs:
            sub_text = sub if isinstance(sub, str) else sub.get("text", "")
            sub_link = "" if isinstance(sub, str) else sub.get("link", "")
            _add_sub_paragraph(tf, sub_text, sub_link, sub_font_size)

        # 空行间隔（条目之间）
        if idx < len(content_items) - 1:
            spacer = tf.add_paragraph()
            spacer.space_after = Pt(0)
            spacer.space_before = Pt(0)
            spacer_run = spacer.add_run()
            spacer_run.text = ""
            spacer_run.font.size = Pt(6)


def create_weekly_report(date_range, this_week_items, next_week_items,
                         long_term_items, output_path):
    """
    生成周报 PPTX 文件

    Args:
        date_range: 日期范围字符串，如 "7.7-7.11"
        this_week_items: 本周工作列表
        next_week_items: 下周计划列表
        long_term_items: 长期任务列表
        output_path: 输出文件路径
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 使用空白布局
    blank_layout = prs.slide_layouts[6]  # 空白布局

    # --- Slide 1: 本周工作 ---
    slide1 = prs.slides.add_slide(blank_layout)
    _add_title_slide(slide1, f"本周工作  {date_range}", this_week_items,
                     is_first_slide=True)

    # --- Slide 2: 下周计划 ---
    slide2 = prs.slides.add_slide(blank_layout)
    _add_title_slide(slide2, f"下周计划  {date_range}", next_week_items)

    # --- Slide 3: 长期任务 ---
    slide3 = prs.slides.add_slide(blank_layout)
    title = "长期任务" if long_term_items else "其他事项"
    _add_title_slide(slide3, title, long_term_items or [{"text": "无", "sub": []}])

    prs.save(output_path)
    return output_path


# ============================================================
# CLI 入口
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="周报 PPT 生成器（支持语雀链接点击跳转）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 命令行传参
  python generate_weekly_ppt.py --date-range "7.7-7.11" \\
      --this-week "PDR算法优化" "台球杆测试" \\
      --next-week "羽毛球项目" "文档整理" \\
      --long-term "触觉传感器" \\
      --output "周报_7.7-7.11.pptx"

  # JSON 文件输入（支持语雀链接）
  python generate_weekly_ppt.py --json data.json --output "周报.pptx"

JSON 格式（link 字段可选）:
  {
    "date_range": "7.7-7.11",
    "this_week": [
      {
        "text": "PDR算法优化",
        "link": "https://wit-motion.yuque.com/sd8v4v/project_doc/za8is2g4u5lure9r",
        "sub": ["完成5组测试", "准确率92%"]
      },
      {
        "text": "台球杆测试",
        "link": "https://wit-motion.yuque.com/sd8v4v/project_doc/dytf4uxpabba8mwc"
      }
    ],
    "next_week": [
      {"text": "羽毛球项目", "sub": ["挥拍速度指标"]}
    ],
    "long_term": [
      {"text": "触觉传感器"}
    ]
  }

链接触发规则:
  - 主条目有 link → 标题行末尾显示 "📎语雀文档"（蓝色下划线可点击）
  - 子条目有 link → 子条目末尾显示 "📎链接"（蓝色下划线可点击）
  - link 字段省略 → 不显示链接
        """
    )

    parser.add_argument("--date-range", "-d", help="日期范围，如 7.7-7.11")
    parser.add_argument("--this-week", nargs="*", default=[],
                        help="本周工作项列表")
    parser.add_argument("--next-week", nargs="*", default=[],
                        help="下周计划项列表")
    parser.add_argument("--long-term", nargs="*", default=[],
                        help="长期任务列表")
    parser.add_argument("--output", "-o", default="周报.pptx",
                        help="输出文件路径 (默认: 周报.pptx)")
    parser.add_argument("--json", "-j", help="从 JSON 文件读取数据")

    args = parser.parse_args()

    # 从 JSON 读取
    if args.json:
        with open(args.json, "r", encoding="utf-8") as f:
            data = json.load(f)
        date_range = data["date_range"]
        this_week = data.get("this_week", [])
        next_week = data.get("next_week", [])
        long_term = data.get("long_term", [])
        output_path = args.output
    else:
        if not args.date_range:
            print("错误: 必须提供 --date-range 或 --json", file=sys.stderr)
            sys.exit(1)
        date_range = args.date_range
        this_week = [{"text": t, "sub": []} for t in args.this_week]
        next_week = [{"text": t, "sub": []} for t in args.next_week]
        long_term = [{"text": t, "sub": []} for t in args.long_term]
        output_path = args.output

    result = create_weekly_report(
        date_range=date_range,
        this_week_items=this_week,
        next_week_items=next_week,
        long_term_items=long_term,
        output_path=output_path,
    )

    print(f"[OK] 周报 PPT 已生成: {result}")


if __name__ == "__main__":
    main()
